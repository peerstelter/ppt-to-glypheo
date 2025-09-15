
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_to_csv_by_color.py (rev3)
CSV: FOLIENNUMMER;DEUTSCH;ENGLISCH

Neu:
- Foliennummern-Platzhalter werden standardmäßig IMMER ignoriert (auch bei --no-skip-placeholders).
- Optionales Einschalten via --include-slide-numbers.
- Defensive Bereinigung: Eine alleinstehende oder am Zeilenende angehängte Foliennummer
  wird aus der Textspalte entfernt.
"""

import argparse
import csv
import sys
import re
from collections import Counter
from typing import List, Optional, Sequence, Tuple

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError as e:
    print("Fehlendes Paket. Bitte installieren Sie:\n  pip install python-pptx pillow", file=sys.stderr)
    raise

RGB = Tuple[int, int, int]

def hex_to_rgb(s: str) -> RGB:
    s = s.strip().lstrip("#")
    if len(s) not in (6, 3):
        raise ValueError(f"Ungültiger Hex-Farbwert: {s}")
    if len(s) == 3:
        s = "".join(ch*2 for ch in s)
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def rgb_to_hex(rgb: RGB) -> str:
    return "{:02X}{:02X}{:02X}".format(*rgb)

def channel_close(a: int, b: int, tol: int) -> bool:
    return abs(a - b) <= tol

def rgb_close(a: RGB, b: RGB, tol: int) -> bool:
    return channel_close(a[0], b[0], tol) and channel_close(a[1], b[1], tol) and channel_close(a[2], b[2], tol)

def get_run_rgb(run) -> Optional[RGB]:
    try:
        fc = run.font.color
    except Exception:
        return None
    if fc is None:
        return None
    try:
        if fc.rgb is not None:
            return (fc.rgb[0], fc.rgb[1], fc.rgb[2])
    except Exception:
        pass
    return None

def _is_specific_placeholder(node, target_names: Sequence[str]) -> bool:
    try:
        if getattr(node, "is_placeholder", False):
            name = str(node.placeholder_format.type)
            return any(key in name for key in target_names)
    except Exception:
        pass
    return False

def iter_shape(shp, skip_placeholders=True, include_slide_numbers=False):
    """Iteriert Shapes, optional Platzhalter überspringen. Foliennummern werden immer übersprungen,
    außer include_slide_numbers=True."""
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shp.shapes:
            yield from iter_shape(s, skip_placeholders=skip_placeholders, include_slide_numbers=include_slide_numbers)
    elif shp.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shp.table.rows:
            for cell in row.cells:
                if cell.text_frame is not None:
                    yield cell
    else:
        # Slide number placeholder?
        if _is_specific_placeholder(shp, ("SLIDE_NUMBER",)) and not include_slide_numbers:
            return
        if skip_placeholders and _is_specific_placeholder(shp, ("DATE", "DATETIME", "FOOTER", "HEADER")):
            return
        yield shp

def extract_runs_with_colors(slide,
                             skip_placeholders: bool = True,
                             include_slide_numbers: bool = False
                             ) -> List[List[Tuple[str, Optional[RGB]]]]:
    """Extrahiert Texte segmentweise entsprechend der Run-Farbe.

    Rückgabe: Liste von Absätzen. Jeder Absatz ist eine Liste aus
    (Textsegment, RGB-Farbe oder None)."""
    paragraphs: List[List[Tuple[str, Optional[RGB]]]] = []
    for shp in slide.shapes:
        for node in iter_shape(shp,
                               skip_placeholders=skip_placeholders,
                               include_slide_numbers=include_slide_numbers):
            tf = getattr(node, "text_frame", None)
            if tf is None:
                continue
            for para in tf.paragraphs:
                if len(para.runs) == 0:
                    continue
                segments: List[Tuple[str, Optional[RGB]]] = []
                buf_txt = ""
                buf_color: Optional[RGB] = None
                for run in para.runs:
                    txt = run.text or ""
                    if not txt:
                        continue
                    color = get_run_rgb(run)
                    if buf_color == color:
                        buf_txt += txt
                    else:
                        if buf_txt:
                            segments.append((buf_txt, buf_color))
                        buf_txt = txt
                        buf_color = color
                if buf_txt:
                    segments.append((buf_txt, buf_color))
                if segments:
                    paragraphs.append(segments)
    return paragraphs

def assign_paragraph_language(run_colors: List[Optional[RGB]],
                              de_colors: Sequence[RGB],
                              en_colors: Sequence[RGB],
                              tol: int,
                              unknown_policy: str) -> str:
    de_count = 0; en_count = 0; unk = 0
    en_defined = len(en_colors) > 0
    for c in run_colors:
        if c is None:
            unk += 1; continue
        if any(rgb_close(c, d, tol) for d in de_colors):
            de_count += 1
        else:
            if en_defined:
                if any(rgb_close(c, e, tol) for e in en_colors):
                    en_count += 1
                else:
                    en_count += 1
            else:
                en_count += 1
    if de_count > en_count:
        return "de"
    if en_count > de_count:
        return "en"
    if unk > 0:
        return "de" if unknown_policy == "german" else ("en" if unknown_policy == "english" else "skip")
    return "de"

def scan_colors(prs: Presentation, skip_placeholders=True, include_slide_numbers=False):
    cnt = Counter()
    for slide in prs.slides:
        for paragraph in extract_runs_with_colors(
            slide,
            skip_placeholders=skip_placeholders,
            include_slide_numbers=include_slide_numbers,
        ):
            for _, c in paragraph:
                if c is not None:
                    cnt[c] += 1
    return cnt

def parse_color_list(s: Optional[str]) -> List[RGB]:
    if not s: return []
    out: List[RGB] = []
    for part in s.split(","):
        part = part.strip()
        if not part: continue
        out.append(hex_to_rgb(part))
    return out

def interactive_mapping(cnt) -> Tuple[List[RGB], List[RGB]]:
    if not cnt:
        print("Keine expliziten RGB-Farben gefunden. Standard wird genutzt (Deutsch=FFFFFF).")
        return [hex_to_rgb("FFFFFF")], []
    print("\nGefundene Textfarben:")
    items = sorted(cnt.items(), key=lambda kv: (-kv[1], rgb_to_hex(kv[0])))
    for i, (rgb, n) in enumerate(items, 1):
        print(f"[{i}] {rgb_to_hex(rgb)}  Vorkommen: {n}")
    de_in = input("Deutsch-Farben (Nummern oder Hex, Komma-getrennt) [Default: FFFFFF]: ").strip()
    de_colors: List[RGB] = [hex_to_rgb("FFFFFF")] if not de_in else []
    if de_in:
        for p in [x.strip() for x in de_in.split(",") if x.strip()]:
            if p.isdigit() and 1 <= int(p) <= len(items):
                de_colors.append(items[int(p)-1][0])
            else:
                de_colors.append(hex_to_rgb(p))
    en_in = input("Englisch-Farben (Nummern oder Hex, optional): ").strip()
    en_colors: List[RGB] = []
    if en_in:
        for p in [x.strip() for x in en_in.split(",") if x.strip()]:
            if p.isdigit() and 1 <= int(p) <= len(items):
                en_colors.append(items[int(p)-1][0])
            else:
                en_colors.append(hex_to_rgb(p))
    return de_colors, en_colors

def _strip_trailing_slide_number(text: str, slide_no: int) -> str:
    """Entfernt eine isolierte oder am Ende angehängte Foliennummer."""
    if not text:
        return text
    # Zeilenweise prüfen
    lines = text.splitlines()
    for i, ln in enumerate(lines):
        # exakte Zahl
        if ln.strip() == str(slide_no):
            lines[i] = ""
            continue
        # am Ende angehängt: "... text 12"
        if ln.rstrip().endswith(" " + str(slide_no)):
            # nur entfernen, wenn vor der Zahl Whitespace ist
            lines[i] = re.sub(rf"\s{slide_no}$", "", ln.rstrip())
    cleaned = "\n".join(l for l in lines if l is not None)
    return cleaned

def extract_to_csv(
    pptx_path: str,
    csv_path: str,
    de_colors: Sequence[RGB],
    en_colors: Sequence[RGB],
    tol: int,
    unknown_policy: str,
    skip_placeholders: bool = True,
    include_slide_numbers: bool = False,
) -> None:
    prs = Presentation(pptx_path)
    rows: List[List[str]] = [["FOLIENNUMMER", "DEUTSCH", "ENGLISCH"]]

    for idx, slide in enumerate(prs.slides, start=1):
        german_parts: List[str] = []
        english_parts: List[str] = []
        paragraphs = extract_runs_with_colors(
            slide,
            skip_placeholders=skip_placeholders,
            include_slide_numbers=include_slide_numbers,
        )
        if not paragraphs:
            rows.append([str(idx), "", ""])
            continue

        for segments in paragraphs:
            para_has_de = False
            para_has_en = False
            for text, color in segments:
                lang = assign_paragraph_language(
                    [color], de_colors, en_colors, tol, unknown_policy
                )
                if lang == "de":
                    german_parts.append(text)
                    para_has_de = True
                elif lang == "en":
                    english_parts.append(text)
                    para_has_en = True
            if para_has_de:
                german_parts.append("\n")
            if para_has_en:
                english_parts.append("\n")

        german = "".join(german_parts).rstrip("\n")
        english = "".join(english_parts).rstrip("\n")

        # Sicherheit: Foliennummer aus Text am Ende entfernen
        german = _strip_trailing_slide_number(german, idx)
        english = _strip_trailing_slide_number(english, idx)

        rows.append([str(idx), german, english])

    with open(csv_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f, delimiter=";", quoting=csv.QUOTE_MINIMAL)
        writer.writerows(rows)

def main(argv: Optional[Sequence[str]] = None) -> int:
    p = argparse.ArgumentParser(description="PPTX → CSV (FOLIENNUMMER;DEUTSCH;ENGLISCH) per Schriftfarbe.")
    p.add_argument("pptx", help="Pfad zur .pptx-Datei")
    p.add_argument("csv", help="Pfad zur Ausgabe-CSV")
    p.add_argument("--de-color", dest="de_color", default=None,
                   help="Deutsch-Farbe(n) als Hex, mehrere mit Komma (Default: FFFFFF)")
    p.add_argument("--en-color", dest="en_color", default=None,
                   help="Englisch-Farbe(n) als Hex, mehrere mit Komma (Default: alle nicht-deutschen)")
    p.add_argument("--tolerance", type=int, default=8,
                   help="Farbtoleranz je Kanal (0..255), Default: 8")
    p.add_argument("--unknown-policy", choices=["german", "english", "skip"], default="german",
                   help="Zuordnung für Runs ohne explizite Farbe. Default: german")
    p.add_argument("--no-skip-placeholders", action="store_true",
                   help="Nicht automatisch Datum/Fußzeile/Kopfzeile überspringen (Foliennummer bleibt ausgeschlossen).")
    p.add_argument("--include-slide-numbers", action="store_true",
                   help="Foliennummern-Platzhalter als Text mit ausgeben.")
    p.add_argument("--interactive", action="store_true",
                   help="Scannt Farben und fragt interaktiv nach Mapping de/en.")
    args = p.parse_args(argv)

    de_colors = parse_color_list(args.de_color) if args.de_color else [hex_to_rgb("FFFFFF")]
    en_colors = parse_color_list(args.en_color)

    prs = Presentation(args.pptx)
    if args.interactive:
        cnt = scan_colors(prs, skip_placeholders=not args.no_skip_placeholders, include_slide_numbers=args.include_slide_numbers)
        de_colors, en_colors = interactive_mapping(cnt)

    extract_to_csv(
        pptx_path=args.pptx,
        csv_path=args.csv,
        de_colors=de_colors,
        en_colors=en_colors,
        tol=args.tolerance,
        unknown_policy=args.unknown_policy,
        skip_placeholders=not args.no_skip_placeholders,
        include_slide_numbers=args.include_slide_numbers,
    )
    print(f"Fertig. CSV geschrieben: {args.csv}")
    return 0

if __name__ == "__main__":
    raise SystemExit(main())
